# -*- coding: utf-8 -*-
# @Author: Qiurui Chen
# @Date: 03/10/2023
# @Last Modified by: Qiurui Chen Contact: rachelchen0831@gmail.com


from pptx import Presentation

presentation_path = "data/DanoneSN.pptx"
presentation = Presentation(presentation_path)

text_list = []
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                parag_list = []
                for run in paragraph.runs:
                    parag_list.append(run.text)
                para_text = " ".join(parag_list)
                para_text = para_text + "."
                text_list.append(para_text)


text_list = [elem for elem in text_list if len(elem) > 1]
text_content = ".".join(text_list)
print(text_list)



from langchain.document_loaders import WebBaseLoader

loader = WebBaseLoader("https://lilianweng.github.io/posts/2023-06-23-agent/")
data = loader.load()

from langchain.text_splitter import RecursiveCharacterTextSplitter

text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=0)
all_splits = text_splitter.split_documents(data)

# texts = text_splitter.create_documents(text_list)
texts = text_splitter.create_documents( [text_content])

from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.llms import OpenAI
from langchain.memory import ConversationSummaryMemory

llm = OpenAI(temperature=0)
memory = ConversationSummaryMemory(llm=llm)
memory.save_context({"input": "hi"},{"output": "whats up"})
memory.save_context({"input": "im working on better docs for chatbots"},{"output": "oh, that sounds like a lot of work"})
memory.save_context({"input": "yes, but it's worth the effort"},{"output": "agreed, good docs are important!"})

vectorstore = Chroma.from_documents(documents=all_splits, embedding=OpenAIEmbeddings())
memory = ConversationSummaryMemory(llm=llm,memory_key="chat_history",return_messages=True)

from langchain.chat_models import ChatOpenAI
from langchain.chains import ConversationalRetrievalChain

llm = ChatOpenAI()
retriever = vectorstore.as_retriever()
qa = ConversationalRetrievalChain.from_llm(llm, retriever=retriever, memory=memory)

qa("How do agents use Task decomposition?")
qa("What are the various ways to implemet memory to support it?")